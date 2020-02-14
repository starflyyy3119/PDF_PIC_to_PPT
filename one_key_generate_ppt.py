import re
import os
import time 
import fitz
import pptx
from pptx.util import Inches
import tkinter as tk
from tkinter import *

def pdf2pic(path, pic_path, imagecount):

    # 从 pdf 中提取图片
    # :param path: pdf 路径
    # :param pic_path: 图片保存的路径
    # :return
    
    t0 = time.process_time()
    # 使用正则表达式来查找图片
    checkXO = r"/Type(?= */XObject)"
    checkIM = r"/Subtype(?= */Image)"
    
    # 打开pdf
    doc = fitz.open(path)
    # 图片计数
    imgcount = imagecount
    lenXREF = doc._getXrefLength()
    
    # 打印pdf信息
    print("文件名:{}, 页数:{}, 对象:{}".format(path, len(doc), lenXREF-1))
    
    # 遍历每一个对象
    for i in range(1, lenXREF):
        # 定义对象字符串
        text = doc._getXrefString(i)
        isXObject = re.search(checkXO, text)
        # 使用正则表达式查看是否是图片
        isImage = re.search(checkIM, text)
        # 如果不是对象也不是图片，则continue
        if not isXObject or not isImage:
            continue
        imgcount = imgcount + 1
        # 根据索引生成图像
        pix = fitz.Pixmap(doc, i)
        # 根据pdf的路径生成图片的名称
        new_name = "img{}.png".format(imgcount)
    
        # 如果pix.n<5,可以直接存为PNG
        if pix.n < 5:
            pix.writePNG(os.path.join(pic_path, new_name))
        # 否则先转换CMYK
        else:
            pix0 = fitz.Pixmap(fitz.csRGB, pix)
            pix0.writePNG(os.path.join(pic_path, new_name))
            pix0 = None
        # 释放资源
        pix = None
        t1 = time.process_time()
        print("运行时间:{}s".format(t1 - t0))
        print("提取了{}张图片".format(imgcount))
        print(new_name)
    return imgcount

def generate_PPT(pic_dir):
    pptFile = pptx.Presentation()
    picFiles = [pic_dir + "\\" + pic for pic in os.listdir(pic_dir)]

    for pic in picFiles:
        slide = pptFile.slides.add_slide(pptFile.slide_layouts[1])
        slide.shapes.add_picture(pic, Inches(0), Inches(0), Inches(10), Inches(7.5))

    pptFile.save(pic_dir + "\\" + '我的ppt.pptx')
def get_file_dir():
    return os.path.join(os.path.expanduser('~'), "Desktop")
class MainWindow(tk.Frame):
    def get_filename(self, event):
        self.file_dir = get_file_dir()
        self.file_name = self.file_name_entry.get()
        self.pic_name = self.pic_name_entry.get()
        self.file_path_dir = self.file_dir + "\\" + self.file_name
        self.pic_path_dir = self.file_dir + "\\" + self.pic_name
    def cancel(self, event):
        exit()
    def get_name(self):
        return self.file_path_dir, self.pic_path_dir
    def __init__(self):
        self.window = Tk()

        self.window.title("一键生成ppt")

        # 输入pdf文件夹名
        self.Label_file = Label(self.window, text = "输入pdf文件夹名")
        self.Label_file.grid(row = 0)
        self.file_name_entry = tk.Entry(self.window)
        self.file_name_entry.grid(row = 0, column = 1)

        # 输入生成图片的文件名
        self.Label_pic = Label(self.window, text = "输入图片文件夹名")
        self.Label_pic.grid(row = 1)
        self.pic_name_entry = tk.Entry(self.window)
        self.pic_name_entry.grid(row = 1, column = 1)

        # 确认按钮
        self.Button_ensure = tk.Button(self.window, text = "确认")
        self.Button_ensure.grid(row = 2)

        # 取消按钮
        self.Button_cancel = tk.Button(self.window, text = "取消")
        self.Button_cancel.grid(row = 2, column = 1)

        # 绑定确定按钮
        self.Button_ensure.bind("<Button-1>", self.get_filename)
        
        # 绑定取消按钮
        self.Button_cancel.bind("<Button-1>", self.cancel)
        self.window.mainloop()
        

if __name__=='__main__':
    window = MainWindow()
    file_path_dir, pic_path_dir = window.get_name()
    imagecount = 0
    paths = [file_path_dir + "\\" + l for l in list(os.listdir(file_path_dir))]

    # 创建保存图片的文件夹
    if os.path.exists(pic_path_dir):
        print("文件夹已存在，请重新创建新文件夹！")
        raise SystemExit
    else:
        os.mkdir(pic_path_dir)
    
    for path in paths:
        imagecount = pdf2pic(path, pic_path_dir, imagecount)
    
    generate_PPT(pic_path_dir)
    